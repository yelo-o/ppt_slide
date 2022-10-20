from pptx import Presentation

# 변수 지정 Presentation() > prs
prs = Presentation()
title_slide_layout = prs.slide_layouts[0] # ppt의 슬라이드 레이아웃 지정

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')