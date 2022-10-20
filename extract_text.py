from pptx import Presentation

# 추출할 ppt 파일 위치 입력
prs = Presentation('C:/Users/mingyu/PycharmProjects/ppt_slide/test.pptx')

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = [] # ppt 텍스트 내용이 리스트 안에 담김

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
print(text_runs) # 내용 출력